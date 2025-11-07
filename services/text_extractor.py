# # text_extractor.py
# from __future__ import annotations
# import os
# import re
# import subprocess
# from pathlib import Path
# from tempfile import TemporaryDirectory
# from typing import Optional

# # ---- OpenAI Whisper（用於音檔/影片語音轉文字） ----
# from openai import OpenAI

# # ---- 文字檔處理相依（保持輕量） ----
# # PDF
# from pdfminer.high_level import extract_text as pdf_extract_text  # type: ignore
# # DOCX
# import docx  # type: ignore
# # PPTX
# from pptx import Presentation  # type: ignore

# from pydub import AudioSegment  # 用來偵測音量
# from services import video_utils  

# # HTML（可選）：若沒裝 bs4 也能退化成簡單正則
# try:
#     from bs4 import BeautifulSoup  # type: ignore
#     _HAS_BS4 = True
# except Exception:
#     _HAS_BS4 = False


# # =========================
# # 基本設定
# # =========================
# # Whisper 模型：官方雲端轉錄
# WHISPER_MODEL = os.getenv("WHISPER_MODEL", "whisper-1")
# _openai_client: Optional[OpenAI] = None

# def _client() -> OpenAI:
#     global _openai_client
#     if _openai_client is None:
#         _openai_client = OpenAI()  # 從環境變數讀 OPENAI_API_KEY
#     return _openai_client


# # 支援的副檔名
# TEXT_EXTS  = {".txt", ".html", ".htm", ".pdf", ".docx", ".pptx"}
# AUDIO_EXTS = {".mp3", ".wav", ".m4a"}
# VIDEO_EXTS = {".mp4", ".mov", ".m4v"}
# IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".bmp"}

# def _is_audio(path: Path) -> bool:
#     return path.suffix.lower() in AUDIO_EXTS

# def _is_video(path: Path) -> bool:
#     return path.suffix.lower() in VIDEO_EXTS

# def _is_text(path: Path) -> bool:
#     return path.suffix.lower() in TEXT_EXTS

# def _is_image(path: Path) -> bool:      
#     return path.suffix.lower() in IMAGE_EXTS


# # =========================
# # 公用：文字清理
# # =========================
# def _normalize_text(s: str) -> str:
#     """基礎清理：統一換行、去 BOM、收斂多餘空白行。"""
#     if not s:
#         return ""
#     # 移除 BOM
#     s = s.replace("\ufeff", "")
#     # 統一換行
#     s = s.replace("\r\n", "\n").replace("\r", "\n")
#     # 收斂連續空白行
#     s = re.sub(r"\n{3,}", "\n\n", s)
#     return s.strip()


# # =========================
# # 純文字 / HTML
# # =========================
# def _read_txt(path: Path) -> str:
#     try:
#         with open(path, "r", encoding="utf-8") as f:
#             return _normalize_text(f.read())
#     except UnicodeDecodeError:
#         # 回退：嘗試 gbk / big5 等（避免爆掉）
#         for enc in ("big5", "gbk", "latin1"):
#             try:
#                 with open(path, "r", encoding=enc, errors="ignore") as f:
#                     return _normalize_text(f.read())
#             except Exception:
#                 pass
#         # 最後一招：binary 解碼
#         with open(path, "rb") as f:
#             return _normalize_text(f.read().decode("utf-8", errors="ignore"))

# def _html_to_text(html: str) -> str:
#     if _HAS_BS4:
#         soup = BeautifulSoup(html, "html.parser")
#         # 移除 script/style
#         for t in soup(["script", "style"]):
#             t.decompose()
#         text = soup.get_text("\n")
#         return _normalize_text(text)
#     # 無 bs4：簡單移除標籤
#     text = re.sub(r"(?is)<script.*?>.*?</script>", "", html)
#     text = re.sub(r"(?is)<style.*?>.*?</style>", "", text)
#     text = re.sub(r"(?s)<[^>]+>", " ", text)
#     return _normalize_text(text)

# def _read_html(path: Path) -> str:
#     raw = _read_txt(path)
#     return _html_to_text(raw)


# # =========================
# # PDF / DOCX / PPTX
# # =========================
# def _read_pdf(path: Path) -> str:
#     return _normalize_text(pdf_extract_text(str(path)) or "")

# def _read_docx(path: Path) -> str:
#     try:
#         d = docx.Document(str(path))
#     except Exception:
#         # 偶見解析問題，改用二進位讀取忽略錯誤
#         d = docx.Document(path)
#     parts = [p.text for p in d.paragraphs if p.text]
#     return _normalize_text("\n".join(parts))

# def _read_pptx(path: Path) -> str:
#     prs = Presentation(str(path))
#     parts = []
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if hasattr(shape, "text") and shape.text:
#                 parts.append(shape.text)
#     return _normalize_text("\n".join(parts))

# def _detect_audio_volume(video_path: str) -> float:
#     """
#     用 ffmpeg 抽取音訊後，用 pydub 偵測平均音量 (dBFS)。
#     回傳值越小代表越安靜，例如 -60 幾乎無聲。
#     """
#     _require_ffmpeg()
#     from tempfile import TemporaryDirectory
#     import subprocess
#     from pathlib import Path

#     with TemporaryDirectory() as td:
#         tmp_audio = Path(td) / "probe.wav"
#         cmd = [
#             "ffmpeg", "-y", "-i", str(video_path),
#             "-vn", "-ac", "1", "-ar", "16000",
#             "-f", "wav", str(tmp_audio)
#         ]
#         subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
#         try:
#             audio = AudioSegment.from_wav(tmp_audio)
#             return audio.dBFS
#         except Exception:
#             return -90.0  # 若無法解析視為極低音量



# # =========================
# # 音檔（Whisper）
# # =========================
# def extract_from_audio(path: str | Path) -> str:
#     """
#     使用 OpenAI Whisper 轉錄音檔（.mp3 / .wav / .m4a）。
#     回傳：轉錄後的全文字串
#     """
#     p = Path(path)
#     with open(p, "rb") as f:
#         # Whisper 雲端 API：以分鐘計價
#         transcript = _client().audio.transcriptions.create(
#             model=WHISPER_MODEL,
#             file=f,
#             response_format="text",  # 直接拿純文字
#             # language="zh",        # 如多為中文可打開；預設自動偵測
#             # temperature=0,
#         )
#     return _normalize_text(transcript)


# # =========================
# # 影片（ffmpeg 抽音 → Whisper）
# # =========================
# def _require_ffmpeg():
#     from shutil import which
#     if which("ffmpeg") is None:
#         raise RuntimeError(
#             "缺少 ffmpeg，請先安裝並加入 PATH（Windows 請安裝 ffmpeg.exe；macOS: brew install ffmpeg）。"
#         )

# def extract_from_video(path: str | Path) -> str:
#     """
#     影片：
#     若有語音 → Whisper + GPT-4o Frame Caption 雙通道融合
#     若無語音 → GPT-4o Frame Caption 單通道摘要
#     """
#     _require_ffmpeg()
#     src = str(path)

#     # --- 新增：音量偵測 ---
#     loudness = _detect_audio_volume(src)
#     print(f"[DEBUG] 平均音量 dBFS = {loudness:.2f}")

#     # 判斷是否有語音（閾值可調）
#     has_audio = loudness > -40

#     if not has_audio:
#         # 🔸 無聲影片：只跑 GPT-4o 畫面摘要
#         print("[INFO] 偵測到無聲影片 → 進行 Frame-based Caption 摘要")
#         captions_text, vision_cost = video_utils.generate_captions(src)
#         return captions_text, vision_cost

#     # 🔸 有聲影片：同時跑 Whisper + Frame Caption + 融合
#     print("[INFO] 偵測到有聲影片 → 啟動雙通道融合模式")

#     # 1. 轉錄語音
#     audio_text = extract_from_video_audioonly(src)

#     # 2. 生成畫面描述
#     captions_text, vision_cost = video_utils.generate_captions(src)

#     # 3. 融合兩者（Whisper + Caption）
#     merged_text = video_utils.fuse_text(audio_text, captions_text)

#     return merged_text, vision_cost


# def extract_from_video_audioonly(src: str) -> str:
#     """
#     專供有聲影片使用的音訊轉文字（保持原 extract_from_video 流程）
#     """
#     with TemporaryDirectory() as td:
#         audio_path = str(Path(td) / "audio.wav")
#         cmd = [
#             "ffmpeg", "-y", "-i", src,
#             "-vn", "-ac", "1", "-ar", "16000", "-f", "wav", audio_path
#         ]
#         ret = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
#         if ret.returncode != 0:
#             raise RuntimeError("ffmpeg 抽取音訊失敗。")
#         return extract_from_audio(audio_path)


# def extract_from_image(path: str | Path):
#     """
#     單張圖片：呼叫 GPT-4o 視覺模型做醫學描述。
#     回傳：描述文字 + 預估 vision_cost
#     """
#     p = Path(path)
#     # 直接呼叫 video_utils 內的單張圖片 caption 函式（下面第 2 部分會加）
#     caption, vision_cost = video_utils.caption_single_image(p)
#     return _normalize_text(caption), vision_cost



# # =========================
# # 統一入口（供外部呼叫）
# # =========================
# def extract_any(path: str | Path) -> str:
#     """
#     依副檔名自動選擇解析方式，回傳 (文字內容, vision_cost)。
#     - 文字檔：TXT / HTML / PDF / DOCX / PPTX
#     - 音檔：MP3 / WAV / M4A（Whisper）
#     - 影片：MP4 / MOV / M4V（ffmpeg 抽音 + Whisper）
#     - 圖片：JPG / PNG / BMP（GPT-4o 視覺摘要）
#     """
#     p = Path(path)
#     ext = p.suffix.lower()

#     # 音檔
#     if _is_audio(p):
#         text = extract_from_audio(p)
#         return text, 0.0

#     # 影片
#     if _is_video(p):
#         return extract_from_video(p)
    
#     # #圖片
#     if _is_image(p):
#         text, vision_cost = extract_from_image(p)
#         return text, vision_cost

#     # 純文字
#     if ext == ".txt":
#         text = _read_txt(p)
#         return text, 0.0

#     # HTML
#     if ext in {".html", ".htm"}:
#         text = _read_html(p)
#         return text, 0.0

#     # PDF / DOCX / PPTX
#     if ext == ".pdf":
#         text = _read_pdf(p)
#         return text, 0.0
#     if ext == ".docx":
#         text = _read_docx(p)
#         return text, 0.0
#     if ext == ".pptx":
#         text = _read_pptx(p)
#         return text, 0.0

#     # 不支援的副檔名
#     return "", 0.0



# text_extractor.py
# 這個模組負責「把各種檔案類型 → 轉成文字」，並且在需要的時候回傳 vision 成本：
# - 純文字 / HTML / PDF / DOCX / PPTX：讀文字
# - 音檔 (mp3/wav/m4a)：用 OpenAI Whisper 轉錄
# - 影片 (mp4/mov/m4v)：偵測有沒有聲音 → 有聲跑 Whisper + 畫面 caption、無聲只跑畫面 caption
# - 圖片 (jpg/png/bmp)：呼叫 GPT-4o 視覺模型做描述
#
# 最外面的入口是：extract_any(path) → 回傳 (文字內容, vision_cost)

from __future__ import annotations
import os
import re
import subprocess
from pathlib import Path
from tempfile import TemporaryDirectory
from typing import Optional

# ---- OpenAI Whisper（用於音檔/影片語音轉文字） ----
from openai import OpenAI

# ---- 文字檔處理相依（保持輕量） ----
# PDF
from pdfminer.high_level import extract_text as pdf_extract_text  # type: ignore
# DOCX
import docx  # type: ignore
# PPTX
from pptx import Presentation  # type: ignore

from pydub import AudioSegment  # 用來偵測音量（dBFS）
from services import video_utils  # 自己的工具，負責影片 frame caption + 圖片 caption + 文本融合等

# HTML（可選）：若沒裝 bs4 也能退化成簡單正則
try:
    from bs4 import BeautifulSoup  # type: ignore
    _HAS_BS4 = True
except Exception:
    _HAS_BS4 = False


# =========================
# 基本設定：Whisper 客戶端
# =========================

# Whisper 模型名稱（從環境變數 WHISPER_MODEL 讀，預設 "whisper-1"）
WHISPER_MODEL = os.getenv("WHISPER_MODEL", "whisper-1")

# OpenAI client（lazy init：第一次用時才建立）
_openai_client: Optional[OpenAI] = None

def _client() -> OpenAI:
    """
    取得共用的 OpenAI client 實例。
    - 若尚未建立，會用預設設定（從 OPENAI_API_KEY 環境變數）建立一個。
    """
    global _openai_client
    if _openai_client is None:
        _openai_client = OpenAI()  # 從環境變數讀 OPENAI_API_KEY
    return _openai_client


# =========================
# 檔案類型判斷
# =========================

# 支援的副檔名分類
TEXT_EXTS  = {".txt", ".html", ".htm", ".pdf", ".docx", ".pptx"}
AUDIO_EXTS = {".mp3", ".wav", ".m4a"}
VIDEO_EXTS = {".mp4", ".mov", ".m4v"}
IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".bmp"}

def _is_audio(path: Path) -> bool:
    """判斷是否為音檔"""
    return path.suffix.lower() in AUDIO_EXTS

def _is_video(path: Path) -> bool:
    """判斷是否為影片檔"""
    return path.suffix.lower() in VIDEO_EXTS

def _is_text(path: Path) -> bool:
    """判斷是否為文字/Office 類檔案"""
    return path.suffix.lower() in TEXT_EXTS

def _is_image(path: Path) -> bool:
    """判斷是否為圖片檔"""
    return path.suffix.lower() in IMAGE_EXTS


# =========================
# 公用：文字清理
# =========================

def _normalize_text(s: str) -> str:
    """
    基礎清理：
    - 去掉 UTF-8 BOM
    - 統一換行符號為 '\n'
    - 把三個以上連續空行縮成兩個
    """
    if not s:
        return ""
    # 移除 BOM
    s = s.replace("\ufeff", "")
    # 統一換行
    s = s.replace("\r\n", "\n").replace("\r", "\n")
    # 收斂連續空白行
    s = re.sub(r"\n{3,}", "\n\n", s)
    return s.strip()


# =========================
# 純文字 / HTML
# =========================

def _read_txt(path: Path) -> str:
    """
    讀取純文字檔：
    - 預設用 UTF-8
    - 若失敗，依序嘗試 big5 / gbk / latin1
    - 最後仍失敗就以 binary 讀取再盡量 decode
    """
    try:
        with open(path, "r", encoding="utf-8") as f:
            return _normalize_text(f.read())
    except UnicodeDecodeError:
        # 回退：嘗試 gbk / big5 等（避免爆掉）
        for enc in ("big5", "gbk", "latin1"):
            try:
                with open(path, "r", encoding=enc, errors="ignore") as f:
                    return _normalize_text(f.read())
            except Exception:
                pass
        # 最後一招：binary 解碼
        with open(path, "rb") as f:
            return _normalize_text(f.read().decode("utf-8", errors="ignore"))

def _html_to_text(html: str) -> str:
    """
    把 HTML 字串轉成純文字：
    - 若有安裝 bs4：用 BeautifulSoup 解析、移除 script/style
    - 否則用簡單正則移除標籤
    """
    if _HAS_BS4:
        soup = BeautifulSoup(html, "html.parser")
        # 移除 script/style
        for t in soup(["script", "style"]):
            t.decompose()
        text = soup.get_text("\n")
        return _normalize_text(text)

    # 無 bs4：簡單移除標籤（效果沒那麼好，但至少不會爆）
    text = re.sub(r"(?is)<script.*?>.*?</script>", "", html)
    text = re.sub(r"(?is)<style.*?>.*?</style>", "", text)
    text = re.sub(r"(?s)<[^>]+>", " ", text)
    return _normalize_text(text)

def _read_html(path: Path) -> str:
    """讀取 HTML 檔，並轉成純文字"""
    raw = _read_txt(path)
    return _html_to_text(raw)


# =========================
# PDF / DOCX / PPTX
# =========================

def _read_pdf(path: Path) -> str:
    """
    使用 pdfminer.high_level.extract_text 直接抽取 PDF 文字。
    注意：這裡沒有做 OCR，如果是掃描 PDF 可能抽不到字。
    """
    return _normalize_text(pdf_extract_text(str(path)) or "")

def _read_docx(path: Path) -> str:
    """
    使用 python-docx 讀取 Word (.docx) 文字。
    若遇解析問題，會改用「傳 Path 物件」的方式再試一次。
    """
    try:
        d = docx.Document(str(path))
    except Exception:
        # 偶見解析問題，改用二進位讀取忽略錯誤
        d = docx.Document(path)
    parts = [p.text for p in d.paragraphs if p.text]
    return _normalize_text("\n".join(parts))

def _read_pptx(path: Path) -> str:
    """
    使用 python-pptx 讀取簡報 (.pptx) 文字：
    - 逐頁 (slide) 逐個 shape
    - 對有 text 屬性的 shape 抽取文字
    """
    prs = Presentation(str(path))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    return _normalize_text("\n".join(parts))


def _detect_audio_volume(video_path: str) -> float:
    """
    用 ffmpeg 抽取音訊後，用 pydub 偵測平均音量 (dBFS)。
    - 回傳值越小代表越安靜，例如 -60 幾乎無聲
    - 這裡用來判斷「影片是否有可用語音」
    """
    _require_ffmpeg()
    from tempfile import TemporaryDirectory
    import subprocess
    from pathlib import Path

    with TemporaryDirectory() as td:
        tmp_audio = Path(td) / "probe.wav"
        cmd = [
            "ffmpeg", "-y", "-i", str(video_path),
            "-vn", "-ac", "1", "-ar", "16000",
            "-f", "wav", str(tmp_audio)
        ]
        subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        try:
            audio = AudioSegment.from_wav(tmp_audio)
            return audio.dBFS
        except Exception:
            # 若無法解析，視為超安靜
            return -90.0  # 若無法解析視為極低音量



# =========================
# 音檔（Whisper）
# =========================

def extract_from_audio(path: str | Path) -> str:
    """
    使用 OpenAI Whisper 轉錄音檔（.mp3 / .wav / .m4a）。
    回傳：轉錄後的全文字串（已做基本清理）
    """
    p = Path(path)
    with open(p, "rb") as f:
        # Whisper 雲端 API：以分鐘計價
        transcript = _client().audio.transcriptions.create(
            model=WHISPER_MODEL,
            file=f,
            response_format="text",  # 直接拿純文字
            # language="zh",        # 如多為中文可打開；預設自動偵測
            # temperature=0,
        )
    # transcript 在 response_format="text" 時，是字串
    return _normalize_text(transcript)


# =========================
# 影片（ffmpeg 抽音 → Whisper）
# =========================

def _require_ffmpeg():
    """
    確認系統中有 ffmpeg 可以用：
    - Windows：需安裝 ffmpeg，並把路徑加入 PATH
    - macOS：可用 brew install ffmpeg
    """
    from shutil import which
    if which("ffmpeg") is None:
        raise RuntimeError(
            "缺少 ffmpeg，請先安裝並加入 PATH（Windows 請安裝 ffmpeg.exe；macOS: brew install ffmpeg）。"
        )

def extract_from_video(path: str | Path):
    """
    影片抽取內容策略：

    1. 先偵測音量（dBFS）判斷影片是否有語音：
       - 若接近無聲（例如 < -40 dBFS）：視為「無聲影片」
       - 若音量 > -40 dBFS：視為「有聲影片」

    2A. 無聲影片：
        - 不跑 Whisper
        - 只跑 GPT-4o Frame-based Caption（看畫面做摘要）
        - 由 video_utils.generate_captions 負責

    2B. 有聲影片：
        - 跑 Whisper 抽音轉文字（extract_from_video_audioonly）
        - 跑 GPT-4o Frame-based Caption（看畫面）
        - 再由 video_utils.fuse_text(audio_text, captions_text) 做「語音 + 畫面」融合

    回傳：
        (merged_text, vision_cost)
        - merged_text：最後融合後的文字內容
        - vision_cost：畫面分析（GPT-4o vision）估算的成本
    """
    _require_ffmpeg()
    src = str(path)

    # --- 新增：音量偵測 ---
    loudness = _detect_audio_volume(src)
    print(f"[DEBUG] 平均音量 dBFS = {loudness:.2f}")

    # 判斷是否有語音（閾值可調整）
    has_audio = loudness > -40

    if not has_audio:
        # 🔸 無聲影片：只跑 GPT-4o 畫面摘要
        print("[INFO] 偵測到無聲影片 → 進行 Frame-based Caption 摘要")
        captions_text, vision_cost = video_utils.generate_captions(src)
        return captions_text, vision_cost

    # 🔸 有聲影片：同時跑 Whisper + Frame Caption + 融合
    print("[INFO] 偵測到有聲影片 → 啟動雙通道融合模式")

    # 1. 轉錄語音（只抽音訊）
    audio_text = extract_from_video_audioonly(src)

    # 2. 生成畫面描述
    captions_text, vision_cost = video_utils.generate_captions(src)

    # 3. 融合兩者（Whisper + Caption）
    merged_text = video_utils.fuse_text(audio_text, captions_text)

    return merged_text, vision_cost


def extract_from_video_audioonly(src: str) -> str:
    """
    專供「有聲影片」使用的音訊轉文字工具：
    - 使用 ffmpeg 從影片中抽出音訊 (wav)
    - 再丟給 extract_from_audio（Whisper）轉文字
    """
    with TemporaryDirectory() as td:
        audio_path = str(Path(td) / "audio.wav")
        cmd = [
            "ffmpeg", "-y", "-i", src,
            "-vn", "-ac", "1", "-ar", "16000", "-f", "wav", audio_path
        ]
        ret = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        if ret.returncode != 0:
            raise RuntimeError("ffmpeg 抽取音訊失敗。")
        return extract_from_audio(audio_path)


def extract_from_image(path: str | Path):
    """
    單張圖片（例如：X 光片、皮膚病灶照片等）：
    - 呼叫 GPT-4o 視覺模型做醫學描述
    - 實際邏輯封裝在 video_utils.caption_single_image

    回傳：
        (描述文字, vision_cost)
    """
    p = Path(path)
    # 直接呼叫 video_utils 內的單張圖片 caption 函式
    caption, vision_cost = video_utils.caption_single_image(p)
    return _normalize_text(caption), vision_cost



# =========================
# 統一入口（供外部呼叫）
# =========================

def extract_any(path: str | Path) -> tuple[str, float]:
    """
    依副檔名自動選擇解析方式，回傳 (文字內容, vision_cost)。

    - 文字檔：TXT / HTML / PDF / DOCX / PPTX
        → 只回文字，vision_cost = 0.0

    - 音檔：MP3 / WAV / M4A（Whisper）
        → 把語音轉成文字，vision_cost = 0.0（因為這是語音，不是視覺）

    - 影片：MP4 / MOV / M4V
        → 可能同時用 Whisper（語音） + GPT-4o vision（畫面摘要）
        → 最終文字 merged_text，vision_cost 由 video_utils 回傳

    - 圖片：JPG / PNG / BMP
        → GPT-4o 視覺摘要，回傳 (文字, vision_cost)
    """
    p = Path(path)
    ext = p.suffix.lower()

    # 音檔
    if _is_audio(p):
        text = extract_from_audio(p)
        return text, 0.0

    # 影片
    if _is_video(p):
        # extract_from_video 已經回傳 (文字, vision_cost)
        return extract_from_video(p)
    
    # 圖片
    if _is_image(p):
        text, vision_cost = extract_from_image(p)
        return text, vision_cost

    # 純文字
    if ext == ".txt":
        text = _read_txt(p)
        return text, 0.0

    # HTML
    if ext in {".html", ".htm"}:
        text = _read_html(p)
        return text, 0.0

    # PDF / DOCX / PPTX
    if ext == ".pdf":
        text = _read_pdf(p)
        return text, 0.0
    if ext == ".docx":
        text = _read_docx(p)
        return text, 0.0
    if ext == ".pptx":
        text = _read_pptx(p)
        return text, 0.0

    # 不支援的副檔名：回傳空字串 + 0 成本
    return "", 0.0
